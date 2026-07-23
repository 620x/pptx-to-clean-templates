"""Юниты движка автолечения — в дополнение к сквозному web-тесту."""

from __future__ import annotations

from pptx import Presentation

from x6pptxgen.heal.engine import _add_layout, _clear_layout


def test_add_layout_creates_valid_part(tmp_path):
    """Низкоуровневое создание layout-части: part + rels + sldLayoutIdLst.

    Путь работает, когда семейств больше, чем свободных layouts, — в обычных
    сценариях не выполняется, поэтому проверяем напрямую.
    """
    prs = Presentation()
    master = prs.slide_masters[0]
    before = len(master.slide_layouts)

    new_layout = _add_layout(master, master.slide_layouts[0])
    _clear_layout(new_layout)
    new_layout._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
    ).set("name", "AUTO test")

    # Новый layout виден мастеру и пригоден для add_slide
    assert len(master.slide_layouts) == before + 1
    slide = prs.slides.add_slide(new_layout)
    assert slide.slide_layout.name == "AUTO test"

    # Файл переживает сохранение и повторное открытие
    out = tmp_path / "with_new_layout.pptx"
    prs.save(out)
    reopened = Presentation(str(out))
    names = [l.name for l in reopened.slide_masters[0].slide_layouts]
    assert "AUTO test" in names
    assert len(reopened.slides._sldIdLst) == 1

    # id в sldLayoutIdLst уникальны
    from pptx.oxml.ns import qn

    lst = reopened.slide_masters[0]._element.find(qn("p:sldLayoutIdLst"))
    ids = [x.get("id") for x in lst]
    assert len(ids) == len(set(ids))


def test_subtitle_role_and_examples_in_prompts(tmp_path):
    """Короткая крупная зона -> подзаголовок idx 11; prompt несёт пример донора."""
    import base64

    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    from x6pptxgen.heal.engine import heal_template

    png = tmp_path / "p.png"
    png.write_bytes(base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
        "AAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="))
    prs = Presentation()
    blank = prs.slide_masters[0].slide_layouts[6]
    for i in range(2):
        s = prs.slides.add_slide(blank)
        t = s.shapes.add_textbox(Emu(500000), Emu(400000), Emu(7000000), Emu(800000))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = f"День {i + 1} в Риме"; r.font.size = Pt(30)
        # короткое крупное название отеля — раньше дропалось
        h = s.shapes.add_textbox(Emu(500000), Emu(3500000), Emu(5000000), Emu(600000))
        hr = h.text_frame.paragraphs[0].add_run()
        hr.text = f"Hotel Bulgari {i + 1}"; hr.font.size = Pt(22)
        b = s.shapes.add_textbox(Emu(500000), Emu(1500000), Emu(7000000), Emu(1800000))
        br = b.text_frame.paragraphs[0].add_run()
        br.text = f"Длинное описание дня номер {i + 1}, " * 5; br.font.size = Pt(13)
    src = tmp_path / "d.pptx"; prs.save(src)

    report = heal_template(src, tmp_path / "h.pptx")
    roles = {p["role"] for p in report["layouts"][0]["placeholders"]}
    assert "subtitle" in roles, "имя отеля должно стать подзаголовком, а не дропом"

    healed = Presentation(str(tmp_path / "h.pptx"))
    prompts = []
    for lay in healed.slide_masters[0].slide_layouts:
        if lay.name.startswith("AUTO"):
            for sp in lay.shapes._spTree.iter(qn("p:sp")):
                ph = sp.find(qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph"))
                if ph is not None:
                    prompts.append("".join(
                        t.text or "" for t in sp.iter(qn("a:t"))))
    assert any("напр.:" in p for p in prompts), "prompt должен нести пример донора"


def test_table_becomes_column_placeholders(tmp_path):
    from pptx.util import Emu, Pt

    from x6pptxgen.heal.engine import heal_template

    prs = Presentation()
    blank = prs.slide_masters[0].slide_layouts[6]
    for i in range(2):
        s = prs.slides.add_slide(blank)
        t = s.shapes.add_textbox(Emu(500000), Emu(400000), Emu(7000000), Emu(700000))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = f"Программа тура {i + 1}"; r.font.size = Pt(28)
        table = s.shapes.add_table(3, 2, Emu(500000), Emu(1500000),
                                   Emu(7000000), Emu(3000000)).table
        for ri in range(3):
            table.cell(ri, 0).text = f"{ri + 1} день"
            table.cell(ri, 1).text = f"описание дня {ri + 1} слайда {i + 1}"
    src = tmp_path / "t.pptx"; prs.save(src)

    report = heal_template(src, tmp_path / "th.pptx")
    idxs = {p["idx"] for p in report["layouts"][0]["placeholders"]}
    assert {20, 21} <= idxs, "таблица должна стать колонками idx 20/21"
