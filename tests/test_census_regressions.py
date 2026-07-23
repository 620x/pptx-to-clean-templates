"""Регрессионные тесты на слепые зоны переписи слайдов (находки ревью Фазы 1).

Каждый тест — минимальная версия воспроизведённого бага:
1. picture placeholder, заполненный картинкой (p:pic сохраняет p:ph), — это
   ИСПОЛЬЗОВАНИЕ шаблона, а не «дизайн на слайде»;
2. фон через p:cSld/p:bg (bgPr blipFill) — третий механизм фото-фона, должен
   считаться грязью на слайдах;
3. картинка внутри сжатой группы — не full-bleed, масштаб ext/chExt учитывается;
4. кривые файлы (нет sldSz, нет связи с темой) — деградация, а не traceback.
"""

from __future__ import annotations

import base64

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu

from x6pptxgen.audit.checks import run_checks
from x6pptxgen.template.inventory import build_inventory

# 1x1 png — python-pptx читает заголовки сам, Pillow не нужен
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
    "AAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


@pytest.fixture()
def tiny_png(tmp_path):
    path = tmp_path / "tiny.png"
    path.write_bytes(_TINY_PNG)
    return path


def _blockers(findings):
    return [f["code"] for f in findings if f["severity"] == "blocker"]


def test_picture_placeholder_counts_as_usage(tmp_path, tiny_png):
    """Фото в picture placeholder не должно давать ложный блокер design-on-slides."""
    prs = Presentation()
    layout = prs.slide_masters[0].slide_layouts[8]  # 'Picture with Caption'
    for _ in range(2):
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                ph.insert_picture(str(tiny_png))
            elif ph.placeholder_format.idx == 0:
                ph.text = "Заголовок"
    path = tmp_path / "pic_ph.pptx"
    prs.save(path)

    inv = build_inventory(path)
    slide_inv = inv["slides"][0]

    kinds = {p["kind"] for p in slide_inv["placeholders_used"]}
    assert "picture" in kinds, "заполненный pic-ph должен попадать в placeholders_used"
    assert all(p["is_placeholder"] for p in slide_inv["pictures"])
    assert "design-on-slides" not in _blockers(run_checks(inv))


def test_slide_background_blipfill_is_detected(tmp_path, tiny_png):
    """Фон «Формат фона > рисунок» на слайде — грязь, которую нельзя пропустить."""
    prs = Presentation()
    blank = prs.slide_masters[0].slide_layouts[6]
    for i in range(2):
        slide = prs.slides.add_slide(blank)
        # Честно получаем image part + rId, затем превращаем в p:bg
        pic = slide.shapes.add_picture(str(tiny_png), Emu(0), Emu(0), Emu(10), Emu(10))
        r_id = pic._element.find(qn("p:blipFill") + "/" + qn("a:blip")).get(qn("r:embed"))
        pic._element.getparent().remove(pic._element)

        csld = slide._element.find(qn("p:cSld"))
        bg = etree.SubElement(csld, qn("p:bg"))
        bgpr = etree.SubElement(bg, qn("p:bgPr"))
        blip_fill = etree.SubElement(bgpr, qn("a:blipFill"))
        etree.SubElement(blip_fill, qn("a:blip")).set(qn("r:embed"), r_id)
        etree.SubElement(bgpr, qn("a:effectLst"))
        csld.insert(0, bg)  # p:bg обязан быть первым ребёнком cSld

        box = slide.shapes.add_textbox(Emu(500000), Emu(500000), Emu(4000000), Emu(800000))
        box.text_frame.text = f"Подпись поверх фона {i + 1}"

    path = tmp_path / "bg_deck.pptx"
    prs.save(path)

    inv = build_inventory(path)
    assert inv["slides"][0]["background_fill"] == "picture"
    assert "design-on-slides" in _blockers(run_checks(inv))


def test_group_scaled_picture_is_not_full_bleed(tmp_path, tiny_png):
    """Картинка «на весь слайд» внутри группы, сжатой до 10%, — не фон."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_picture(
        str(tiny_png), Emu(0), Emu(0), prs.slide_width, prs.slide_height
    )
    # Сжимаем группу до 10%, оставляя «детскую» систему координат полноразмерной —
    # ровно это пишет PowerPoint при масштабировании группы.
    xfrm = group._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
    ext = xfrm.find(qn("a:ext"))
    chext = xfrm.find(qn("a:chExt"))
    chext.set("cx", str(prs.slide_width))
    chext.set("cy", str(prs.slide_height))
    ext.set("cx", str(prs.slide_width // 10))
    ext.set("cy", str(prs.slide_height // 10))

    path = tmp_path / "grouped.pptx"
    prs.save(path)

    inv = build_inventory(path)
    pic = inv["slides"][0]["pictures"][0]
    assert pic["in_group"]
    assert not pic["full_bleed"]
    assert pic["pct_w"] <= 15


def test_degenerate_files_degrade_instead_of_crashing(tmp_path, default_pptx):
    """Нет p:sldSz / нет связи с темой — находка аудита, а не traceback."""
    # без p:sldSz
    prs = Presentation(str(default_pptx))
    el = prs.part._element.find(qn("p:sldSz"))
    el.getparent().remove(el)
    no_size = tmp_path / "no_size.pptx"
    prs.save(no_size)
    inv = build_inventory(no_size)
    assert inv["slide_size"]["assumed_default"]
    assert inv["problems"]

    # без связи мастера с темой
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    prs = Presentation(str(default_pptx))
    master_part = prs.slide_masters[0].part
    theme_rels = [rid for rid, rel in master_part.rels.items() if rel.reltype == RT.THEME]
    for rid in theme_rels:
        master_part.drop_rel(rid)
    no_theme = tmp_path / "no_theme.pptx"
    prs.save(no_theme)
    inv = build_inventory(no_theme)
    assert inv["fonts"]["themes_by_master"]["m0"] == {"major_latin": "", "minor_latin": ""}
    assert any("тем" in p for p in inv["problems"])

    findings = run_checks(inv)
    assert any(f["code"] == "structure-degraded" for f in findings)
