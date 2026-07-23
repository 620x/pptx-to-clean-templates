"""Инвентаризация шаблона: мастера, layouts, плейсхолдеры, слайды, шрифты.

Почему здесь много сырого XML при живом python-pptx: библиотека прячет ровно то,
что важно аудитору (проверено, см. docs/research/):
- её словарь placeholders при дублях idx молча возвращает первый матч;
- обход shapes не спускается в группы p:grpSp;
- для встроенных шрифтов API нет вообще;
- эффективная типографика не резолвится (font.name/size == None значит «унаследовано»).
Поэтому структуру берём из python-pptx, а всё «подозрительное» перечитываем
по lxml-дереву напрямую. Доступ к приватным _element / _spTree — общепринятый
приём: публичного XML-API у python-pptx нет.

Важно: плейсхолдер на слайде — не всегда p:sp. Заполненный picture placeholder
становится p:pic, таблица — p:graphicFrame, и оба СОХРАНЯЮТ p:ph (проверено
ревью Фазы 1). Перепись обязана видеть все три вида, иначе чистый шаблон
получает ложный блокер.
"""

from __future__ import annotations

import hashlib
from pathlib import Path

from pptx.oxml.ns import qn

from x6pptxgen.template.opener import load_template

INVENTORY_VERSION = "0.2"

# Размер слайда по умолчанию (16:9), если p:sldSz отсутствует — файл кривой,
# но аудит должен выдать находку, а не traceback.
_DEFAULT_SLIDE_W, _DEFAULT_SLIDE_H = 12192000, 6858000

# Наследование геометрии layout -> master идёт по ТИПУ плейсхолдера через
# фиксированную таблицу (правило PowerPoint, зеркалит python-pptx; проверено
# в docs/research/ooxml-dirty.json). idx между layout и master не участвует.
_LAYOUT_TYPE_TO_MASTER_TYPE = {
    "body": "body",
    "chart": "body",
    "clipArt": "body",
    "dgm": "body",
    "media": "body",
    "obj": "body",
    "pic": "body",
    "subTitle": "body",
    "tbl": "body",
    "ctrTitle": "title",
    "title": "title",
    "dt": "dt",
    "ftr": "ftr",
    "sldNum": "sldNum",
}

# Типы, которые add_slide не клонирует на новый слайд — для генерации это «фон»,
# а не контент (проверено: жёсткий список в python-pptx).
NON_CONTENT_PH_TYPES = {"dt", "ftr", "sldNum"}

# Шейпы, способные нести p:ph, и их nv-контейнеры (у каждого вида свой).
_NVPR_BY_TAG = {}


def _init_tag_tables() -> None:
    # qn() требует загруженных namespace-таблиц python-pptx; инициализируем лениво,
    # чтобы модуль импортировался без сюрпризов.
    if _NVPR_BY_TAG:
        return
    _NVPR_BY_TAG[qn("p:sp")] = "p:nvSpPr"
    _NVPR_BY_TAG[qn("p:pic")] = "p:nvPicPr"
    _NVPR_BY_TAG[qn("p:graphicFrame")] = "p:nvGraphicFramePr"


# Куда macOS кладёт шрифты — для эвристики «найдём ли шрифт для замера».
_FONT_DIRS = [
    Path("/System/Library/Fonts"),
    Path("/System/Library/Fonts/Supplemental"),
    Path("/Library/Fonts"),
    Path.home() / "Library/Fonts",
]


def build_inventory(path: str | Path) -> dict:
    """Собрать полный инвентарь шаблона в виде простого dict (сериализуется в JSON).

    Никогда не бросает исключений на «кривых, но открывающихся» файлах:
    деградации записываются в inventory["problems"] и становятся находками аудита.
    """
    _init_tag_tables()
    path = Path(path)
    prs = load_template(path)
    problems: list[str] = []

    inventory: dict = {
        "inventory_version": INVENTORY_VERSION,
        "source_file": path.name,
        "file_sha1": hashlib.sha1(path.read_bytes()).hexdigest(),
        "slide_size": _slide_size(prs, problems),
        "masters": [],
        "slides": [],
        "fonts": _fonts_info(prs, problems),
        "problems": problems,
    }

    layout_usage: dict[str, list[int]] = {}
    for n, slide in enumerate(prs.slides, start=1):
        layout_part = str(slide.slide_layout.part.partname)
        layout_usage.setdefault(layout_part, []).append(n)

    for mi, master in enumerate(prs.slide_masters):
        master_geometry = _placeholder_geometry_by_type(master.shapes._spTree)
        m: dict = {
            "id": f"m{mi}",
            "name": master.name or f"master {mi}",
            "layouts": [],
        }
        for li, layout in enumerate(master.slide_layouts):
            part = str(layout.part.partname)
            m["layouts"].append({
                "id": f"m{mi}/l{li}",
                "part": part.rsplit("/", 1)[-1],
                "name": layout.name,
                "placeholders": _layout_placeholders(layout, master_geometry),
                "pictures": _pictures(layout.shapes._spTree, inventory["slide_size"]),
                "background_fill": _background_fill(layout._element),
                "used_by_slides": layout_usage.get(part, []),
            })
        inventory["masters"].append(m)

    for n, slide in enumerate(prs.slides, start=1):
        inventory["slides"].append(_slide_census(n, slide, inventory))

    return inventory


# --- размер слайда -----------------------------------------------------------

def _slide_size(prs, problems: list[str]) -> dict:
    cx, cy = prs.slide_width, prs.slide_height
    assumed = cx is None or cy is None
    if assumed:
        # Легально-опциональный p:sldSz отсутствует — файл странный, но не смертельно.
        problems.append(
            "В presentation.xml нет p:sldSz — размер слайда принят 16:9 по умолчанию, "
            "проценты картинок и будущие замеры неточны."
        )
        cx, cy = _DEFAULT_SLIDE_W, _DEFAULT_SLIDE_H
    cx, cy = int(cx), int(cy)
    return {
        "cx_emu": cx,
        "cy_emu": cy,
        "inches": f"{cx / 914400:.2f} x {cy / 914400:.2f}",
        "orientation": "portrait" if cy > cx else "landscape",
        "assumed_default": assumed,
    }


# --- плейсхолдеры ------------------------------------------------------------

def _ph_of(el):
    """Вернуть p:ph шейпа (p:sp / p:pic / p:graphicFrame) или None."""
    nv_tag = _NVPR_BY_TAG.get(el.tag)
    if nv_tag is None:
        return None
    return el.find(qn(nv_tag) + "/" + qn("p:nvPr") + "/" + qn("p:ph"))


def _ph_type_idx(ph) -> tuple[str, int]:
    # По схеме OOXML отсутствующий type == "obj", отсутствующий idx == 0 —
    # это легальные дефолты, а не ошибка шаблона.
    return ph.get("type") or "obj", int(ph.get("idx") or 0)


def _xfrm(el) -> dict | None:
    """Собственная геометрия шейпа. У graphicFrame p:xfrm лежит прямо в нём,
    у sp/pic — внутри p:spPr."""
    if el.tag == qn("p:graphicFrame"):
        xfrm = el.find(qn("p:xfrm"))
    else:
        sppr = el.find(qn("p:spPr"))
        xfrm = sppr.find(qn("a:xfrm")) if sppr is not None else None
    if xfrm is None:
        return None
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return None
    return {
        "x_emu": int(off.get("x")),
        "y_emu": int(off.get("y")),
        "cx_emu": int(ext.get("cx")),
        "cy_emu": int(ext.get("cy")),
    }


def _shape_name(el) -> str:
    nv_tag = _NVPR_BY_TAG.get(el.tag)
    if nv_tag is None:
        return ""
    cnv = el.find(qn(nv_tag) + "/" + qn("p:cNvPr"))
    return cnv.get("name", "") if cnv is not None else ""


def _iter_ph_capable(sp_tree):
    """Все шейпы поддерева, способные нести p:ph (обход рекурсивный, включая группы)."""
    return sp_tree.iter(qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"))


def _placeholder_geometry_by_type(sp_tree) -> dict[str, dict]:
    """Карта type -> геометрия для плейсхолдеров мастера (источник наследования)."""
    result: dict[str, dict] = {}
    for sp in _iter_ph_capable(sp_tree):
        ph = _ph_of(sp)
        if ph is None:
            continue
        ph_type, _ = _ph_type_idx(ph)
        geo = _xfrm(sp)
        if geo is not None and ph_type not in result:
            result[ph_type] = geo
    return result


def _layout_placeholders(layout, master_geometry: dict[str, dict]) -> list[dict]:
    """Плейсхолдеры layout по сырому spTree — включая дубли idx, которые API прячет."""
    placeholders = []
    seen_idx: dict[int, int] = {}
    for sp in _iter_ph_capable(layout.shapes._spTree):
        ph = _ph_of(sp)
        if ph is None:
            continue
        ph_type, idx = _ph_type_idx(ph)
        seen_idx[idx] = seen_idx.get(idx, 0) + 1

        geo = _xfrm(sp)
        source = "layout"
        if geo is None:
            # Своей геометрии нет — наследуется с мастера по типу.
            master_type = _LAYOUT_TYPE_TO_MASTER_TYPE.get(ph_type, "body")
            geo = master_geometry.get(master_type)
            source = "master" if geo is not None else "missing"

        placeholders.append({
            "idx": idx,
            "type": ph_type,
            "name": _shape_name(sp),
            "is_content": ph_type not in NON_CONTENT_PH_TYPES,
            "geometry": geo,
            "geometry_source": source,
            # Эффективный шрифт/кегль и ёмкость появятся со style-резолвером
            # (Вечера 3-6 плана); честно помечаем, что пока не знаем.
            "effective_font": "unknown",
            "capacity_chars": "unknown",
        })

    duplicated = {idx for idx, count in seen_idx.items() if count > 1}
    for p in placeholders:
        p["duplicate_idx"] = p["idx"] in duplicated
    return placeholders


# --- слайды ------------------------------------------------------------------

def _text_of(sp) -> str:
    body = sp.find(qn("p:txBody"))
    if body is None:
        return ""
    return "".join(t.text or "" for t in body.iter(qn("a:t")))


def _in_group(el) -> bool:
    parent = el.getparent()
    while parent is not None:
        if parent.tag == qn("p:grpSp"):
            return True
        parent = parent.getparent()
    return False


def _group_scale(el) -> tuple[float, float]:
    """Суммарный масштаб от вложенных групп: ext/chExt по каждому предку-группе.

    Группа рисует детей в «детской» системе координат (chOff/chExt) и растягивает
    её в свой ext — картинка «на весь слайд» внутри сжатой группы на деле маленькая.
    """
    sx = sy = 1.0
    parent = el.getparent()
    while parent is not None:
        if parent.tag == qn("p:grpSp"):
            xfrm = parent.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
            if xfrm is not None:
                ext, chext = xfrm.find(qn("a:ext")), xfrm.find(qn("a:chExt"))
                if ext is not None and chext is not None:
                    chx, chy = int(chext.get("cx") or 0), int(chext.get("cy") or 0)
                    if chx and chy:
                        sx *= int(ext.get("cx")) / chx
                        sy *= int(ext.get("cy")) / chy
        parent = parent.getparent()
    return sx, sy


def _background_fill(sld_like) -> str | None:
    """Тип заливки фона p:cSld/p:bg: 'picture' | 'fill' | 'theme-ref' | None.

    'picture' на СЛАЙДЕ — третий механизм фото-фона («Формат фона > рисунок»),
    невидимый через перечень шейпов; на layout/мастере — наоборот, хорошо:
    новые слайды унаследуют его автоматически.
    """
    bg = sld_like.find(qn("p:cSld") + "/" + qn("p:bg"))
    if bg is None:
        return None
    bgpr = bg.find(qn("p:bgPr"))
    if bgpr is not None:
        if bgpr.find(qn("a:blipFill")) is not None:
            return "picture"
        return "fill"
    if bg.find(qn("p:bgRef")) is not None:
        return "theme-ref"
    return None


def _pictures(sp_tree, slide_size: dict, ph_geometry_by_idx: dict[int, dict] | None = None) -> list[dict]:
    """Все p:pic поддерева (обход рекурсивный). Проценты — от размера слайда.

    Нюансы (все из ревью Фазы 1):
    - размер внутри групп корректируется на масштаб группы (ext/chExt);
    - p:pic может быть заполненным picture placeholder'ом БЕЗ своей геометрии —
      тогда берём геометрию одноимённого idx из layout (ph_geometry_by_idx);
    - если геометрию не разрешить, картинку всё равно показываем (pct = None),
      а full_bleed не утверждаем.
    """
    pics = []
    for pic in sp_tree.iter(qn("p:pic")):
        ph = _ph_of(pic)
        geo = _xfrm(pic)
        if geo is None and ph is not None and ph_geometry_by_idx:
            _, idx = _ph_type_idx(ph)
            geo = ph_geometry_by_idx.get(idx)

        entry = {
            "is_placeholder": ph is not None,
            "in_group": _in_group(pic),
            "pct_w": None,
            "pct_h": None,
            "full_bleed": False,
        }
        if geo is not None:
            sx, sy = _group_scale(pic)
            pct_w = geo["cx_emu"] * sx / slide_size["cx_emu"]
            pct_h = geo["cy_emu"] * sy / slide_size["cy_emu"]
            entry.update({
                "pct_w": round(pct_w * 100),
                "pct_h": round(pct_h * 100),
                "full_bleed": pct_w >= 0.9 and pct_h >= 0.9,
            })
        pics.append(entry)
    return pics


def _slide_census(n: int, slide, inventory: dict) -> dict:
    """Перепись одного слайда-примера: плейсхолдеры vs плавающие боксы vs картинки.

    Плейсхолдеры ищем во всех трёх видах шейпов (sp/pic/graphicFrame): заполненный
    картинкой или таблицей плейсхолдер — это ИСПОЛЬЗОВАНИЕ шаблона по назначению,
    а не «дизайн на слайде».
    """
    sp_tree = slide.shapes._spTree

    placeholders_used = []
    floating_boxes = []
    for el in _iter_ph_capable(sp_tree):
        ph = _ph_of(el)
        if ph is not None:
            ph_type, idx = _ph_type_idx(ph)
            if el.tag == qn("p:sp"):
                kind, has_content = "text", bool(_text_of(el).strip())
            elif el.tag == qn("p:pic"):
                kind, has_content = "picture", True
            else:
                kind, has_content = "frame", True
            placeholders_used.append(
                {"idx": idx, "type": ph_type, "kind": kind, "has_content": has_content}
            )
        elif el.tag == qn("p:sp") and _text_of(el).strip():
            # Текст без p:ph — «дизайн на слайде», главный антипаттерн.
            floating_boxes.append({
                "name": _shape_name(el),
                "chars": len(_text_of(el)),
                "in_group": _in_group(el),
            })

    layout_part = str(slide.slide_layout.part.partname).rsplit("/", 1)[-1]
    layout_rec = next(
        (l for m in inventory["masters"] for l in m["layouts"] if l["part"] == layout_part),
        None,
    )
    layout_id = layout_rec["id"] if layout_rec else layout_part
    ph_geometry = (
        {p["idx"]: p["geometry"] for p in layout_rec["placeholders"] if p["geometry"]}
        if layout_rec else {}
    )

    run_fonts = sorted({
        latin.get("typeface")
        for latin in sp_tree.iter(qn("a:latin"))
        if latin.get("typeface")
    })

    return {
        "n": n,
        "layout_id": layout_id,
        "placeholders_used": placeholders_used,
        "floating_text_boxes": floating_boxes,
        "pictures": _pictures(sp_tree, inventory["slide_size"], ph_geometry),
        "background_fill": _background_fill(slide._element),
        "run_fonts": run_fonts,
    }


# --- шрифты ------------------------------------------------------------------

def _fonts_info(prs, problems: list[str]) -> dict:
    # Тема — у КАЖДОГО мастера своя; сверка «шрифты слайдов vs тема» должна идти
    # против темы мастера конкретного слайда, не против slide_masters[0].
    themes: dict[str, dict] = {}
    for mi, master in enumerate(prs.slide_masters):
        themes[f"m{mi}"] = _theme_fonts_of_master(master, f"m{mi}", problems)

    embedded = _embedded_fonts(prs)
    usage = _font_usage(prs)

    referenced = set(usage["master_layouts"]) | set(usage["slides"])
    for theme in themes.values():
        referenced |= {theme["major_latin"], theme["minor_latin"]}
    referenced.discard("")
    availability = {
        face: _system_font_hit(face)
        for face in sorted(f for f in referenced if not f.startswith("+"))
    }

    return {
        "theme": themes.get("m0", {"major_latin": "", "minor_latin": ""}),
        "themes_by_master": themes,
        "embedded": embedded,
        "usage": usage,
        # Эвристика по именам файлов в системных папках шрифтов; точный резолвинг
        # семейства -> файла появится в модуле measure (Вечер 6).
        "system_availability_heuristic": availability,
    }


def _theme_fonts_of_master(master, master_id: str, problems: list[str]) -> dict:
    # У python-pptx нет объектной модели темы — часть theme загружается как
    # generic Part, поэтому парсим её blob напрямую. Отсутствие темы — не повод
    # ронять аудит: фиксируем проблему и продолжаем.
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    try:
        theme_part = master.part.part_related_by(RT.THEME)
    except KeyError:
        problems.append(
            f"У мастера {master_id} нет связи с темой (битые rels) — "
            "шрифты темы для него неизвестны."
        )
        return {"major_latin": "", "minor_latin": ""}

    root = etree.fromstring(theme_part.blob)
    scheme = root.find(".//" + qn("a:fontScheme"))

    def latin_of(tag: str) -> str:
        el = scheme.find(qn(f"a:{tag}") + "/" + qn("a:latin")) if scheme is not None else None
        return el.get("typeface", "") if el is not None else ""

    return {"major_latin": latin_of("majorFont"), "minor_latin": latin_of("minorFont")}


def _embedded_fonts(prs) -> list[dict]:
    result = []
    font_list = prs.part._element.find(qn("p:embeddedFontLst"))
    if font_list is None:
        return result
    for ef in font_list.findall(qn("p:embeddedFont")):
        font = ef.find(qn("p:font"))
        styles = [
            style for style in ("regular", "bold", "italic", "boldItalic")
            if ef.find(qn(f"p:{style}")) is not None
        ]
        result.append({
            "typeface": font.get("typeface", "?") if font is not None else "?",
            "styles": styles,
        })
    return result


def _font_usage(prs) -> dict:
    """Гистограммы литеральных a:latin по зонам: каркас (мастер+layouts) vs слайды.

    Хардкод шрифтов в каркасе вместо ссылок на тему (+mj-lt/+mn-lt) — признак
    «нечестной» темы; шрифты на слайдах, отличные от темы, — признак того, что
    реальный стиль живёт не там, где его будет искать генератор.
    """
    def histogram(elements) -> dict[str, int]:
        counts: dict[str, int] = {}
        for el in elements:
            for latin in el.iter(qn("a:latin")):
                face = latin.get("typeface")
                if face:
                    counts[face] = counts.get(face, 0) + 1
        return dict(sorted(counts.items(), key=lambda kv: -kv[1]))

    scaffolding = [m._element for m in prs.slide_masters]
    scaffolding += [l._element for m in prs.slide_masters for l in m.slide_layouts]
    slides = [s._element for s in prs.slides]

    scaffold_hist = histogram(scaffolding)
    literal = sum(c for face, c in scaffold_hist.items() if not face.startswith("+"))
    total = sum(scaffold_hist.values())

    return {
        "master_layouts": scaffold_hist,
        "slides": histogram(slides),
        "scaffolding_literal_ratio": round(literal / total, 2) if total else 0.0,
    }


def _system_font_hit(typeface: str) -> bool:
    """Грубая эвристика: есть ли файл шрифта с похожим именем в системных папках."""
    needle = typeface.lower().replace(" ", "")
    for font_dir in _FONT_DIRS:
        if not font_dir.is_dir():
            continue
        for f in font_dir.iterdir():
            if needle in f.name.lower().replace(" ", ""):
                return True
    return False
