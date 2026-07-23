"""Открытие файла шаблона (.pptx / .potx / .ppsx).

Зачем этот модуль: python-pptx отказывается открывать .potx — у шаблона другой
content-type главной части, и Presentation() бросает ValueError. Лечится на уровне
zip: перед открытием подменяем content-type части /ppt/presentation.xml на
"…presentation.main+xml" (рецепт проверен в docs/research/python-pptx.json).

Исходный файл пользователя никогда не изменяется — вся работа идёт с копией в памяти.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

# Content-types главной части, которые python-pptx не принимает,
# и тот единственный, который принимает.
_TEMPLATE_CT = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_SLIDESHOW_CT = "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml"
_PRESENTATION_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def load_template(path: str | Path) -> PresentationType:
    """Открыть файл шаблона как Presentation, не трогая файл на диске.

    .pptx открывается как есть; .potx/.ppsx — через патч content-type в копии.
    """
    data = Path(path).read_bytes()

    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        content_types = zf.read("[Content_Types].xml").decode("utf-8")

    if _TEMPLATE_CT not in content_types and _SLIDESHOW_CT not in content_types:
        return Presentation(io.BytesIO(data))

    # Пересобираем zip с исправленным [Content_Types].xml.
    # Подмена строки безопасна: template/slideshow content-type встречается
    # только в Override главной части.
    patched = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(data)) as src,
        zipfile.ZipFile(patched, "w", zipfile.ZIP_DEFLATED) as dst,
    ):
        for item in src.infolist():
            payload = src.read(item.filename)
            if item.filename == "[Content_Types].xml":
                text = payload.decode("utf-8")
                text = text.replace(_TEMPLATE_CT, _PRESENTATION_CT)
                text = text.replace(_SLIDESHOW_CT, _PRESENTATION_CT)
                payload = text.encode("utf-8")
            dst.writestr(item, payload)
    patched.seek(0)
    return Presentation(patched)
