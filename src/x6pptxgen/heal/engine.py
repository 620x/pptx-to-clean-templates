"""Автоматическое лечение шаблона: перенос дизайна со слайдов в layouts.

Обобщение ручного лечения реального корпоративного шаблона в движок без
привязки к конкретному файлу. Конвейер:

1. Кластеризация слайдов-примеров по структурной сигнатуре (копипаст-вёрстка
   даёт повторяющиеся семейства — это и есть будущие layouts).
2. В каждом семействе выбирается слайд-донор, его шейпы получают роли:
   - декор (картинки, градиенты, повторяющийся между слайдами текст) —
     копируется статикой с перепривязкой image/hyperlink-связей;
   - контентные зоны (заголовок/кикер/боди/примечания — по кеглю, позиции
     и длине текста) — становятся плейсхолдерами, типографика переносится
     в lstStyle, чтобы наследоваться «по построению».
3. Семейства записываются в неиспользуемые layouts; когда те кончаются —
   создаются новые layout-части (низкоуровневый OPC).

Автолечение — ЧЕРНОВИК онбординга: результат показывается пользователю
пробой (по слайду на layout) до какой-либо генерации.
"""

from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from x6pptxgen.template.opener import load_template

# Роли контентных зон -> (тип плейсхолдера, idx, prompt). Схема idx едина
# с исходным ручным лечением — меню/генерация ничего о происхождении не знают.
_ROLE_PH = {
    "title": ("title", 0, "Заголовок"),
    "kicker": ("body", 10, "Надзаголовок"),
    "subtitle": ("body", 11, "Подзаголовок"),
    "body": ("body", 1, "Основной текст"),
    "body2": ("body", 2, "Дополнительный текст"),
    "body3": ("body", 3, "Дополнительный текст 2"),
    "footnote": ("body", 13, "Примечания"),
}

MAX_LAYOUTS = 8
MIN_TITLE_PT = 15.0
FOOTNOTE_PT = 9.5


# --- анализ ------------------------------------------------------------------

@dataclass
class TextZone:
    element: object
    chars: int
    size_pt: float
    y_frac: float
    h_frac: float
    text: str
    role: str | None = None


@dataclass
class FamilyPlan:
    signature: tuple
    slide_numbers: list[int]
    donor_index: int          # индекс слайда-донора (0-based)
    name: str = ""
    zones: list[TextZone] = field(default_factory=list)
    static_ids: set = field(default_factory=set)


def _slide_signature(slide_inv: dict) -> tuple:
    """Та же сигнатура, что в аудите (checks._slide_signature)."""
    from x6pptxgen.audit.checks import _slide_signature as sig

    return sig(slide_inv)


def _shape_max_size_pt(sp) -> float:
    sizes = [
        int(rpr.get("sz")) / 100
        for rpr in sp.iter(qn("a:rPr"), qn("a:defRPr"))
        if rpr.get("sz")
    ]
    return max(sizes, default=14.0)


def _text_of(sp) -> str:
    body = sp.find(qn("p:txBody"))
    if body is None:
        return ""
    return "".join(t.text or "" for t in body.iter(qn("a:t")))


def _norm_text(text: str) -> str:
    return " ".join(text.split()).lower()


def _is_placeholder(sp) -> bool:
    return sp.find(qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph")) is not None


def _iter_text_shapes(slide):
    """Текст-несущие p:sp ТОЛЬКО верхнего уровня (не плейсхолдеры).

    В группы не спускаемся сознательно: плейсхолдер внутри p:grpSp — нелегальный
    PML, а зона из группы всё равно не сматчится при записи layout (находка
    ревью: текст в группе съедал роль заголовка и вмораживался декором).
    Группа целиком — декор.
    """
    for sp in slide.shapes._spTree:
        if sp.tag != qn("p:sp") or _is_placeholder(sp):
            continue
        text = _text_of(sp)
        if text.strip():
            yield sp, text


def _set_shape_name(el, name: str) -> None:
    for nv_tag in ("p:nvSpPr", "p:nvPicPr", "p:nvGraphicFramePr"):
        cnv = el.find(qn(nv_tag) + "/" + qn("p:cNvPr"))
        if cnv is not None:
            cnv.set("name", name)
            return


def _xfrm_area(el) -> int | None:
    sppr = el.find(qn("p:spPr"))
    xfrm = sppr.find(qn("a:xfrm")) if sppr is not None else None
    ext = xfrm.find(qn("a:ext")) if xfrm is not None else None
    if ext is None:
        return None
    return int(ext.get("cx")) * int(ext.get("cy"))


def _shape_geometry_frac(sp, slide_w: int, slide_h: int) -> tuple[float, float]:
    sppr = sp.find(qn("p:spPr"))
    xfrm = sppr.find(qn("a:xfrm")) if sppr is not None else None
    if xfrm is None:
        return 0.0, 0.0
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return 0.0, 0.0
    return int(off.get("y")) / slide_h, int(ext.get("cy")) / slide_h


def _assign_roles(zones: list[TextZone]) -> None:
    """Эвристики ролей. Работают по кеглю/позиции/длине — не по содержимому."""
    free = [z for z in zones if z.role is None]
    if not free:
        return

    # Заголовок: максимальный кегль (при равенстве — выше на слайде).
    # Если явных кеглей в деке нет (всё наследуется от мастера и вернулся
    # дефолт), принимаем СТРОГО крупнейшую короткую зону — иначе дек без
    # явной типографики оставался без заголовков вовсе (находка ревью).
    title = max(free, key=lambda z: (z.size_pt, -z.y_frac))
    others = [z for z in free if z is not title]
    strictly_largest = all(z.size_pt < title.size_pt for z in others)
    if title.size_pt >= MIN_TITLE_PT or (
        others and strictly_largest and title.chars <= 120
    ):
        title.role = "title"
    elif not others and title.chars <= 120:
        title.role = "title"

    # Кикер: короткая строка НАД заголовком с кеглем меньше титульного
    if title.role == "title":
        kickers = [
            z for z in free
            if z.role is None and z.y_frac < title.y_frac
            and z.chars <= 60 and z.size_pt < title.size_pt
        ]
        if kickers:
            min(kickers, key=lambda z: title.y_frac - z.y_frac).role = "kicker"

    # Примечания: мелкий кегль в нижней части
    for z in free:
        if z.role is None and z.size_pt <= FOOTNOTE_PT and z.y_frac + z.h_frac > 0.75:
            z.role = "footnote"
            break

    # Основной текст: самый длинный из оставшихся; дальше — дополнительные
    rest = sorted(
        (z for z in free if z.role is None and z.chars >= 30),
        key=lambda z: -z.chars,
    )
    for z, role in zip(rest, ("body", "body2", "body3")):
        z.role = role

    # Короткие КРУПНЫЕ зоны — подзаголовки (имя отеля, площадки...), а не мусор:
    # дроп такой зоны лишал модель слота под название (находка на плотном реальном деке)
    for z in free:
        if z.role is None and z.size_pt >= 18 and z.chars <= 40:
            z.role = "subtitle"
            break
    # Остальное (мелкие короткие подписи) остаётся декором


def _family_name(index: int, slide_inv: dict) -> str:
    pics = slide_inv["pictures"]
    if slide_inv.get("background_fill") == "picture" or any(
        p["full_bleed"] for p in pics
    ):
        hint = "фото-фон"
    elif any((p["pct_w"] or 0) >= 60 for p in pics):
        hint = "фото-полоса"
    elif any((p["pct_h"] or 0) >= 60 for p in pics):
        hint = "фото сбоку"
    elif pics:
        hint = "текст и фото"
    else:
        hint = "текст"
    return f"AUTO {index}: {hint}"


def plan_healing(prs, inventory: dict) -> list[FamilyPlan]:
    """Разбить слайды на семейства и разметить роли зон доноров."""
    slides = list(prs.slides)
    # p:sldSz опционален — тот же фолбэк, что и в инвентаре (находка ревью:
    # аудит такой файл принимал, а лечение падало TypeError'ом)
    slide_w = int(prs.slide_width or 12192000)
    slide_h = int(prs.slide_height or 6858000)

    # Константы уровня ДЕКА: текст, дословно повторяющийся на разных слайдах
    # (шапки, кнопки «ОФИЦИАЛЬНЫЙ САЙТ»...), — всегда декор, в любой семье.
    deck_counts: dict[str, int] = {}
    for slide in slides:
        for _, text in _iter_text_shapes(slide):
            key = _norm_text(text)
            deck_counts[key] = deck_counts.get(key, 0) + 1
    deck_consts = {t for t, c in deck_counts.items() if c > 1}

    families: dict[tuple, FamilyPlan] = {}
    for slide_inv in inventory["slides"]:
        sig = _slide_signature(slide_inv)
        fam = families.setdefault(
            sig, FamilyPlan(signature=sig, slide_numbers=[], donor_index=-1)
        )
        fam.slide_numbers.append(slide_inv["n"])

    plans = sorted(families.values(), key=lambda f: (-len(f.slide_numbers), f.slide_numbers[0]))
    plans = plans[:MAX_LAYOUTS]

    for i, fam in enumerate(plans, 1):
        # Донор — слайд с наибольшим числом текст-зон (самый полный пример)
        candidates = [
            (n, sum(1 for _ in _iter_text_shapes(slides[n - 1])))
            for n in fam.slide_numbers
        ]
        donor_n = max(candidates, key=lambda c: c[1])[0]
        fam.donor_index = donor_n - 1
        fam.name = _family_name(i, inventory["slides"][donor_n - 1])

        donor = slides[fam.donor_index]
        zones: list[TextZone] = []
        for sp, text in _iter_text_shapes(donor):
            y, h = _shape_geometry_frac(sp, slide_w, slide_h)
            zones.append(TextZone(
                element=sp, chars=len(text.strip()),
                size_pt=_shape_max_size_pt(sp), y_frac=y, h_frac=h, text=text,
            ))
        fam.zones = zones

        # Кандидат в заголовки освобождается от const-фильтра: семейство из
        # дословно одинаковых слайдов иначе давало макет без плейсхолдеров.
        title_candidate = max(zones, key=lambda z: (z.size_pt, -z.y_frac), default=None)
        for zone in zones:
            if zone is title_candidate:
                continue
            if _norm_text(zone.text) in deck_consts or zone.chars < 3:
                fam.static_ids.add(id(zone.element))
        _assign_roles([z for z in zones if id(z.element) not in fam.static_ids])

    return plans


# --- перенос в layouts (общие механики ручного лечения) ----------------------

def _copy_with_refs(el, src_part, dst_part):
    """Deepcopy шейпа с перепривязкой r:id-ссылок (картинки + гиперссылки)."""
    el2 = deepcopy(el)
    for blip in el2.iter(qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        if rid:
            image_part = src_part.rels[rid].target_part
            blip.set(qn("r:embed"), dst_part.relate_to(image_part, RT.IMAGE))
    for tag in ("a:hlinkClick", "a:hlinkMouseOver"):
        for hlink in list(el2.iter(qn(tag))):
            rid = hlink.get(qn("r:id"))
            if not rid:
                continue
            rel = src_part.rels.get(rid)
            if rel is not None and rel.is_external:
                hlink.set(
                    qn("r:id"),
                    dst_part.relate_to(rel.target_ref, rel.reltype, is_external=True),
                )
            else:
                hlink.getparent().remove(hlink)
    return el2


def _style_key(rpr) -> tuple:
    """Ключ стиля рана: (кегль, жирность, цвет, шрифт) — для кластеризации."""
    if rpr is None:
        return (None, None, None, None)
    color = None
    clr = rpr.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
    if clr is not None:
        color = clr.get("val")
    latin = rpr.find(qn("a:latin"))
    return (rpr.get("sz"), rpr.get("b"), color,
            latin.get("typeface") if latin is not None else None)


def _paragraph_clusters(txbody) -> list[tuple]:
    """Кластеры стилей абзацев зоны: [(chars, style_key, pPr, rPr), ...].

    Реальные блоки серий разностильные ВНУТРИ (бордовый подзаголовок дня,
    жирное время, обычный текст) — плоский перенос одного стиля давал
    «стену текста» (на деке с плотными разностильными блоками). Кластеризуем по доминирующему
    рану абзаца, порядок — по первому появлению."""
    clusters: dict[tuple, list] = {}
    order: list[tuple] = []
    for p_el in txbody.findall(qn("a:p")):
        best_rpr, best_len, chars = None, -1, 0
        for r_el in p_el.findall(qn("a:r")):
            t = r_el.find(qn("a:t"))
            length = len(t.text or "") if t is not None else 0
            chars += length
            if length > best_len:
                best_len = length
                best_rpr = r_el.find(qn("a:rPr"))
        if not chars:
            continue
        key = _style_key(best_rpr)
        if key not in clusters:
            clusters[key] = [0, key, p_el.find(qn("a:pPr")), best_rpr]
            order.append(key)
        clusters[key][0] += chars
    return [tuple(clusters[k]) for k in order]


def _to_placeholder(sp, ph_type: str, idx: int, prompt: str) -> None:
    """Шейп -> плейсхолдер с переносом типографики в lstStyle.

    Многостильные зоны переносятся УРОВНЯМИ абзацев (родной механизм OOXML):
    lvl1 — базовый стиль (максимум символов), lvl2..lvl4 — остальные стили
    в порядке появления. Модель/пользователь выбирают уровень абзаца,
    рендер наследует его стиль из макета."""
    nvpr = sp.find(qn("p:nvSpPr") + "/" + qn("p:nvPr"))
    ph = etree.SubElement(nvpr, qn("p:ph"))
    if ph_type != "obj":
        ph.set("type", ph_type)
    if idx:
        ph.set("idx", str(idx))
    nvpr.insert(0, ph)

    cnvsppr = sp.find(qn("p:nvSpPr") + "/" + qn("p:cNvSpPr"))
    if cnvsppr is not None and cnvsppr.find(qn("a:spLocks")) is None:
        etree.SubElement(cnvsppr, qn("a:spLocks")).set("noGrp", "1")

    txbody = sp.find(qn("p:txBody"))
    clusters = _paragraph_clusters(txbody)
    if clusters:
        base = max(clusters, key=lambda c: c[0])
        ordered = [base] + [c for c in clusters if c is not base][:3]
    else:
        ordered = []

    lst = txbody.find(qn("a:lstStyle"))
    if lst is None:
        lst = txbody.makeelement(qn("a:lstStyle"), {})
        txbody.insert(list(txbody).index(txbody.find(qn("a:bodyPr"))) + 1, lst)
    for level, (_, _, src_ppr, src_rpr) in enumerate(ordered, 1):
        lvl = etree.SubElement(lst, qn(f"a:lvl{level}pPr"))
        if src_ppr is not None:
            for attr in ("algn", "marL", "indent"):
                if src_ppr.get(attr):
                    lvl.set(attr, src_ppr.get(attr))
            for child in src_ppr:
                if not child.tag.endswith("}defRPr"):
                    lvl.append(deepcopy(child))
        if lvl.find(qn("a:buNone")) is None and src_ppr is None:
            etree.SubElement(lvl, qn("a:buNone"))
        if src_rpr is not None:
            defrpr = deepcopy(src_rpr)
            defrpr.tag = qn("a:defRPr")
            for bad in (defrpr.findall(qn("a:hlinkClick"))
                        + defrpr.findall(qn("a:hlinkMouseOver"))):
                defrpr.remove(bad)
            lvl.append(defrpr)

    for p in txbody.findall(qn("a:p")):
        txbody.remove(p)
    p = etree.SubElement(txbody, qn("a:p"))
    r = etree.SubElement(p, qn("a:r"))
    etree.SubElement(r, qn("a:t")).text = prompt


def _copy_slide_background(donor_slide, layout) -> None:
    """Перенести <p:bg> донорского слайда в cSld макета (с перепривязкой
    картинок фона, если это blipFill). p:bg обязан быть первым в cSld."""
    bg = donor_slide._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if bg is None:
        return
    csld = layout._element.find(qn("p:cSld"))
    old = csld.find(qn("p:bg"))
    if old is not None:
        csld.remove(old)
    bg2 = deepcopy(bg)
    for blip in bg2.iter(qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        if rid:
            image_part = donor_slide.part.rels[rid].target_part
            blip.set(qn("r:embed"), layout.part.relate_to(image_part, RT.IMAGE))
    csld.insert(0, bg2)


def _clear_layout(layout) -> None:
    sp_tree = layout.shapes._spTree
    keep = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in list(sp_tree):
        if child.tag not in keep:
            sp_tree.remove(child)


def _add_layout(master, base_layout):
    """Создать НОВУЮ layout-часть, склонировав base (низкоуровневый OPC:
    part + связи master<->layout + запись в sldLayoutIdLst)."""
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.packuri import PackURI
    from pptx.parts.slide import SlideLayoutPart

    package = base_layout.part.package
    existing = {str(part.partname) for part in package.iter_parts()}
    i = 1
    while f"/ppt/slideLayouts/slideLayout{i}.xml" in existing:
        i += 1
    partname = PackURI(f"/ppt/slideLayouts/slideLayout{i}.xml")

    new_el = deepcopy(base_layout._element)
    part = SlideLayoutPart(partname, CT.PML_SLIDE_LAYOUT, package, new_el)
    r_id = master.part.relate_to(part, RT.SLIDE_LAYOUT)
    part.relate_to(master.part, RT.SLIDE_MASTER)

    lst = master._element.find(qn("p:sldLayoutIdLst"))
    # id обязан быть уникален в документе и >= 2147483648
    max_id = max((int(x.get("id")) for x in lst), default=2147483648)
    entry = etree.SubElement(lst, qn("p:sldLayoutId"))
    entry.set("id", str(max_id + 1))
    entry.set(qn("r:id"), r_id)

    layout = part.slide_layout
    _clear_layout(layout)
    return layout


def _table_to_columns(frame, layout, warnings: list[str], fam_name: str) -> list[dict]:
    """Таблица донора -> до двух текстовых колонок-плейсхолдеров (idx 20, 21).

    Геометрия — из p:xfrm таблицы и реальных ширин a:gridCol; стиль — из первой
    непустой ячейки колонки; prompt — с примером содержимого. Механика
    выверена на ручном лечении реальной программы тура.
    """
    xfrm = frame.find(qn("p:xfrm"))
    tbl = frame.find(".//" + qn("a:tbl"))
    if xfrm is None or tbl is None:
        warnings.append(f"{fam_name}: graphicFrame без таблицы пропущен")
        return []
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        warnings.append(f"{fam_name}: таблица без геометрии пропущена")
        return []

    grid_cols = tbl.findall(".//" + qn("a:gridCol"))
    rows = tbl.findall(qn("a:tr"))
    n_cols = min(len(grid_cols), 2)
    if not n_cols or not rows:
        warnings.append(f"{fam_name}: пустая таблица пропущена")
        return []

    x, y = int(off.get("x")), int(off.get("y"))
    cy = int(ext.get("cy"))

    made = []
    col_x = x
    for ci in range(n_cols):
        col_w = int(grid_cols[ci].get("w"))
        # Пример и стиль — из первой непустой ячейки колонки
        example, sz, color, bold = "", "1100", "000000", False
        for tr in rows:
            cells = tr.findall(qn("a:tc"))
            if ci >= len(cells):
                continue
            text = " ".join(
                t.text or "" for t in cells[ci].iter(qn("a:t"))).strip()
            if text and not example:
                example = text[:60]
                rpr = cells[ci].find(".//" + qn("a:rPr"))
                if rpr is not None:
                    sz = rpr.get("sz") or sz
                    bold = rpr.get("b") == "1"
                    clr = rpr.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
                    if clr is not None:
                        color = clr.get("val")
        idx = 20 + ci
        prompt = (f"Колонка {ci + 1} — напр.: «{example}»" if example
                  else f"Колонка {ci + 1}")
        bold_attr = ' b="1"' if bold else ""
        sp_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{2000 + idx}" name="AUTO TABLE COL {ci + 1}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{col_x}" y="{y}"/><a:ext cx="{col_w}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr lIns="0" tIns="0" rIns="91425" bIns="0" anchor="t" wrap="square"/>
    <a:lstStyle>
      <a:lvl1pPr algn="l">
        <a:lnSpc><a:spcPct val="100000"/></a:lnSpc>
        <a:spcBef><a:spcPts val="600"/></a:spcBef>
        <a:buNone/>
        <a:defRPr sz="{sz}"{bold_attr}><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:r><a:t>{prompt}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
        layout.shapes._spTree.append(etree.fromstring(sp_xml))
        made.append({"role": f"table-col-{ci + 1}", "idx": idx, "size_pt": int(sz) / 100})
        col_x += col_w
    return made


# --- сам heal ----------------------------------------------------------------

def heal_template(template_path: str | Path, out_path: str | Path) -> dict:
    """Вылечить шаблон автоматически. Возвращает отчёт для UI."""
    from x6pptxgen.template.inventory import build_inventory

    template_path, out_path = Path(template_path), Path(out_path)
    inventory = build_inventory(template_path)
    prs = load_template(template_path)
    slides = list(prs.slides)
    master = prs.slide_masters[0]
    slide_w = int(prs.slide_width or 12192000)
    slide_h = int(prs.slide_height or 6858000)

    plans = plan_healing(prs, inventory)
    if not plans:
        raise RuntimeError("В файле нет слайдов-примеров — лечить нечего")

    # Неиспользуемые layouts — цели для перезаписи; кончатся — создадим новые
    used_parts = {str(s.slide_layout.part.partname) for s in slides}
    free_layouts = [
        l for l in master.slide_layouts
        if str(l.part.partname) not in used_parts
    ]

    report_layouts = []
    warnings: list[str] = []
    for fam in plans:
        donor = slides[fam.donor_index]

        # Семейство без единой контентной роли (например, донор-таблица) —
        # не пишем: макет без плейсхолдеров бесполезен для генерации,
        # а его почти пустая проба сбивает пользователя (находка ревью).
        if not any(z.role in _ROLE_PH for z in fam.zones):
            warnings.append(
                f"{fam.name} (слайды {fam.slide_numbers}): не удалось выделить "
                "контентные зоны — семейство пропущено, нужна ручная доводка"
            )
            continue

        if free_layouts:
            layout = free_layouts.pop(0)
            _clear_layout(layout)
        else:
            layout = _add_layout(master, master.slide_layouts[0])

        # Фон донорского слайда (кремовая/цветная заливка p:bg) — в макет:
        # иначе сгенерированные слайды получают белый фон мастера
        # (замечено на реальном деке с цветным фоном)
        _copy_slide_background(donor, layout)

        placeholders = []
        dropped: list[str] = []
        photos = []
        photo_seen = False
        zones_by_id = {id(z.element): z for z in fam.zones}
        slide_area = slide_w * slide_h
        for el in donor.shapes._spTree:
            if el.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
                continue
            if el.tag == qn("p:graphicFrame"):
                # Таблица донора -> текстовые колонки-плейсхолдеры: таблицы вне
                # MVP, но их СОДЕРЖИМОЕ — контент (программа тура!), и терять
                # его молча нельзя (на реальном деке центральная страница выходила пустой)
                cols = _table_to_columns(el, layout, warnings, fam.name)
                placeholders.extend(cols)
                continue
            zone = zones_by_id.get(id(el))
            if zone is not None and zone.role is None and id(el) not in fam.static_ids:
                # Уникальный текст без роли (имя отеля, дата конкретного тура) —
                # НЕ вмораживаем в макет: это контент донора, а не декор бренда
                # (находка ревью: имя отеля из донора протаскивалось во все будущие деки)
                dropped.append(zone.text.strip()[:40])
                continue
            el2 = _copy_with_refs(el, donor.part, layout.part)

            becomes_placeholder = zone is not None and zone.role in _ROLE_PH

            # Крупные фото (не логотипы) — адресуемые фото-зоны: пользователь
            # сможет заменить их своей картинкой per-slide. Помечаем именем
            # AUTO PHOTO n; ВЕСЬ декор после фото по z-порядку (градиенты,
            # логотипы, константные подписи) — AUTO OVERLAY n: при замене фото
            # генератор скопирует его на слайд поверх новой картинки, иначе
            # декор оказался бы закрыт ею.
            if el2.tag == qn("p:pic") and (
                (_xfrm_area(el2) or 0) >= slide_area * 0.10
            ):
                n = len(photos) + 1
                _set_shape_name(el2, f"AUTO PHOTO {n}")
                photos.append({"n": n})
                photo_seen = True
            elif photo_seen and not becomes_placeholder:
                _set_shape_name(el2, f"AUTO OVERLAY {len(photos)}")

            if becomes_placeholder:
                ph_type, idx, prompt = _ROLE_PH[zone.role]
                # Пример донорского содержимого в prompt — главная подсказка
                # модели, ЧТО пишется в эту зону (роль «Основной текст» немая:
                # без примеров получалась каша день/отель)
                example = " ".join(zone.text.split())[:70]
                _to_placeholder(el2, ph_type, idx,
                                f"{prompt} — напр.: «{example}»")
                placeholders.append(
                    {"role": zone.role, "idx": idx, "size_pt": zone.size_pt}
                )
            layout.shapes._spTree.append(el2)
        if dropped:
            warnings.append(
                f"{fam.name}: контент донора не перенесён в макет "
                f"(это правильно): {', '.join(dropped[:5])}"
            )

        layout._element.find(qn("p:cSld")).set("name", fam.name)
        report_layouts.append({
            "name": fam.name,
            "from_slides": fam.slide_numbers,
            "placeholders": placeholders,
            "photo_zones": len(photos),
        })

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return {
        "out": str(out_path),
        "layouts": report_layouts,
        "families": len(plans),
        "slides_total": len(slides),
        "warnings": warnings,
        "menu_prefix": "AUTO",
    }


_PROBE_TEXTS = {
    # Короткие тексты: у доноров боксы порой перекрываются, и длинная проба
    # превращалась в кашу на обложках (находка ревью)
    0: "Проба",
    10: "КИКЕР",
    1: "Текст пробы: типографика и цвета наследуются от макета.",
    2: "Блок 2",
    3: "Блок 3",
    13: "*Примечание",
}


def build_probe(healed_path: str | Path, out_path: str | Path,
                name_prefix: str = "AUTO") -> int:
    """Проба лечения: по слайду на каждый новый layout с ролевыми текстами.

    Это то, что пользователь видит ДО генерации: «вот как будут выглядеть
    слайды из вылеченного шаблона». Стиль нигде не задаётся — только
    наследование, как и в настоящей генерации.
    """
    prs = load_template(healed_path)
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)

    count = 0
    for layout in prs.slide_masters[0].slide_layouts:
        if not layout.name.startswith(name_prefix):
            continue
        slide = prs.slides.add_slide(layout)
        for ph in list(slide.placeholders):
            text = _PROBE_TEXTS.get(ph.placeholder_format.idx)
            if text:
                ph.text = text
            else:
                ph._element.getparent().remove(ph._element)
        count += 1

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return count
